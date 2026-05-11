"""PowerPoint parser — Step 2 of the parsing pipeline.

Extracts a list of chunk dicts from a PPTX file, where each chunk follows
the schema defined in docs/parsing-pipeline.md:

    {
        "source":        str,   # file path or URL
        "page_or_slide": int,   # 1-indexed slide number
        "type":          str,   # "heading" | "body" | "notes" | "transcript"
        "text":          str,   # normalised text content
    }

Compared to the PDF parser, PPTX files carry richer structure: every slide
has a distinct title placeholder, body shapes, and an optional speaker-notes
pane. Speaker notes are often the most valuable content because lecturers
write detailed explanations there that never appear in the slide deck itself.

Public API:
    - parse_pptx(path) : extract title / body / notes chunks per slide

Run directly to parse a file and print its chunk count:

    python -m parsers.pptx_parser path/to/file.pptx
    python python/parsers/pptx_parser.py path/to/file.pptx
"""

from __future__ import annotations

import sys
from typing import List, TypedDict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class Chunk(TypedDict):
    source: str
    page_or_slide: int
    type: str
    text: str


# python-pptx placeholder indices: 0 == title, 1 == body. Anything else is a
# subtitle, footer, content-block, etc., which we treat as body text.
_TITLE_PLACEHOLDER_IDX = 0

# Drop body chunks shorter than this — typically slide labels, page numbers,
# or stray bullet glyphs that survived the text-frame strip.
_BODY_MIN_CHARS = 10

# Speaker notes only become a chunk when they're substantial. The threshold
# is a touch higher than for body text because notes tend to either be empty
# placeholders ("Click to add notes") or full paragraphs of explanation.
_NOTES_MIN_CHARS = 20


def parse_pptx(path: str) -> List[Chunk]:
    """Extract structured chunks from a PPTX deck.

    Emits up to one heading chunk, N body chunks, and one notes chunk per
    slide, in slide order. Picture shapes are skipped — they have no text
    frame to extract from, even when they sometimes pretend to.
    """
    prs = Presentation(path)
    chunks: List[Chunk] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title_text = ""
        body_texts: List[str] = []
        notes_text = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            if (
                shape.is_placeholder
                and shape.placeholder_format.idx == _TITLE_PLACEHOLDER_IDX
            ):
                title_text = text
            else:
                body_texts.append(text)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        if title_text:
            chunks.append(
                {
                    "source": path,
                    "page_or_slide": slide_num,
                    "type": "heading",
                    "text": title_text,
                }
            )

        for body in body_texts:
            if len(body) >= _BODY_MIN_CHARS:
                chunks.append(
                    {
                        "source": path,
                        "page_or_slide": slide_num,
                        "type": "body",
                        "text": body,
                    }
                )

        if len(notes_text) >= _NOTES_MIN_CHARS:
            chunks.append(
                {
                    "source": path,
                    "page_or_slide": slide_num,
                    "type": "notes",
                    "text": notes_text,
                }
            )

    return chunks


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python pptx_parser.py <path-to-pptx>", file=sys.stderr)
        sys.exit(2)

    pptx_path = sys.argv[1]
    parsed = parse_pptx(pptx_path)
    print(len(parsed))
